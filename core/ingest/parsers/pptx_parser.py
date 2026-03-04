from __future__ import annotations

from typing import Dict, List

from core.utils.optional_deps import require


def _shape_text(shape) -> str:
    try:
        value = str(getattr(shape, "text", "") or "").strip()
        return value
    except (OSError, ValueError, TypeError, KeyError, AttributeError, RuntimeError):
        return ""


def parse_pptx(path: str) -> Dict:
    ok, msg = require(
        "pptx",
        "pip install python-pptx",
        "PPTX ingest",
    )
    if not ok:
        return {"text": "", "slides": 0, "notes_found": 0, "error": msg}

    try:
        from pptx import Presentation
    except (OSError, ValueError, TypeError, KeyError, AttributeError, RuntimeError, ImportError, ModuleNotFoundError) as exc:
        return {"text": "", "slides": 0, "notes_found": 0, "error": str(exc)}

    try:
        prs = Presentation(path)
        lines: List[str] = []
        notes_found = 0

        for idx, slide in enumerate(prs.slides, start=1):
            slide_lines: List[str] = []
            for shape in slide.shapes:
                text = _shape_text(shape)
                if text:
                    slide_lines.append(text)

            note_text = ""
            try:
                notes = getattr(slide, "notes_slide", None)
                if notes is not None:
                    note_text = str(getattr(notes.notes_text_frame, "text", "") or "").strip()
            except (OSError, ValueError, TypeError, KeyError, AttributeError, RuntimeError):
                note_text = ""

            if note_text:
                notes_found += 1
                slide_lines.append(f"[notes] {note_text}")

            if slide_lines:
                lines.append(f"[slide {idx}]")
                lines.extend(slide_lines)

        return {
            "text": "\n".join(lines).strip(),
            "slides": len(prs.slides),
            "notes_found": notes_found,
            "error": None,
        }
    except (OSError, ValueError, TypeError, KeyError, AttributeError, RuntimeError) as exc:
        return {"text": "", "slides": 0, "notes_found": 0, "error": str(exc)}
